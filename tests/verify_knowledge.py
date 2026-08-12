from argparse import ArgumentParser
from importlib.util import module_from_spec, spec_from_file_location
from pathlib import Path
import json
import re
import sys
import tempfile

ROOT = Path(__file__).resolve().parents[1]
TOOLS = ROOT / "src/tools/knowledge"


def load_module(name: str, path: Path):
    if str(path.parent) not in sys.path:
        sys.path.insert(0, str(path.parent))
    spec = spec_from_file_location(name, path)
    module = module_from_spec(spec)
    sys.modules[name] = module
    spec.loader.exec_module(module)
    return module


def expect_error(fn, text: str) -> None:
    try:
        fn()
    except (KeyError, RuntimeError, ValueError) as exc:
        assert text.lower() in str(exc).lower(), exc
    else:
        raise AssertionError(f"expected error containing {text!r}")


def layer_1() -> None:
    state = json.loads((ROOT / "feature-list.json").read_text(encoding="utf-8"))
    assert state["allowed_states"] == ["not_started", "active", "blocked", "passing"]
    active = [feature["id"] for feature in state["features"] if feature["state"] == "active"]
    assert active == ["H006"], active
    h006 = next(feature for feature in state["features"] if feature["id"] == "H006")
    assert h006["depends_on"] == ["H005"] and h006["evidence"] is None
    assert h006["verification"] == {
        "layer_1": "python tests/verify_knowledge.py --layer 1",
        "layer_2": "python tests/verify_knowledge.py --layer 2",
        "layer_3": "fresh Hermes process, approved Azure resources, and approved Telegram chat complete docs/plan/Rag.md End-to-End Release Checklist A-G",
    }
    required = (
        TOOLS / "contracts.py", TOOLS / "manifest.py", TOOLS / "azure.py",
        TOOLS / "extract.py", TOOLS / "ingest.py", ROOT / "src/.env.example",
    )
    assert all(path.is_file() for path in required), "knowledge contract files missing"
    env_names = {line.split("=", 1)[0] for line in (ROOT / "src/.env.example").read_text(encoding="utf-8").splitlines() if "=" in line}
    assert env_names == {
        "AZURE_STORAGE_CONNECTION_STRING", "AZURE_STORAGE_CONTAINER", "AZURE_SEARCH_ENDPOINT",
        "AZURE_SEARCH_ADMIN_KEY", "AZURE_SEARCH_QUERY_KEY", "AZURE_SEARCH_INDEX",
        "AZURE_DOCUMENT_INTELLIGENCE_ENDPOINT", "AZURE_DOCUMENT_INTELLIGENCE_KEY",
        "AZURE_OPENAI_ENDPOINT", "AZURE_OPENAI_API_KEY", "AZURE_OPENAI_EMBEDDING_DEPLOYMENT",
    }
    ignored = (ROOT / ".gitignore").read_text(encoding="utf-8")
    assert "src/.runtime/" in ignored and ".env" in ignored
    tracked_text = "\n".join(path.read_text(encoding="utf-8", errors="ignore") for path in ROOT.rglob("*") if path.is_file() and ".git" not in path.parts and ".worktrees" not in path.parts)
    secret_patterns = (r"(?i)api[_-]?key\s*=\s*['\"][^'\"]+", r"(?i)authorization:\s*bearer\s+\S+", r"DefaultEndpointsProtocol=https;AccountName=")
    assert not any(re.search(pattern, tracked_text) for pattern in secret_patterns), "possible committed secret"


def azure_config():
    return {
        "AZURE_STORAGE_CONNECTION_STRING": "DefaultEndpointsProtocol=https;AccountName=test;AccountKey=storage-secret",
        "AZURE_STORAGE_CONTAINER": "knowledge",
        "AZURE_SEARCH_ENDPOINT": "https://search.example",
        "AZURE_SEARCH_ADMIN_KEY": "admin-secret",
        "AZURE_SEARCH_QUERY_KEY": "query-secret",
        "AZURE_SEARCH_INDEX": "knowledge",
        "AZURE_DOCUMENT_INTELLIGENCE_ENDPOINT": "https://document.example",
        "AZURE_DOCUMENT_INTELLIGENCE_KEY": "document-secret",
        "AZURE_OPENAI_ENDPOINT": "https://openai.example",
        "AZURE_OPENAI_API_KEY": "openai-secret",
        "AZURE_OPENAI_EMBEDDING_DEPLOYMENT": "text-embedding-3-small",
    }


def layer_2() -> None:
    contracts = load_module("contracts", TOOLS / "contracts.py")
    manifest_module = load_module("manifest", TOOLS / "manifest.py")
    azure = load_module("azure", TOOLS / "azure.py")
    extraction = load_module("extract", TOOLS / "extract.py")
    ingestion = load_module("ingest", TOOLS / "ingest.py")

    for bad in ("", "../x.pdf", "a/../x.pdf", "C:/x.pdf", "x.exe", "a\\x.pdf"):
        expect_error(lambda value=bad: contracts.validate_source_path(value), "source" if bad != "x.exe" else "unsupported")

    evidence = contracts.Evidence(
        chunk_id="doc:g1:c1", content="Refunds require approval.", source="Policy.pdf",
        source_path="policies/Policy.pdf", page_number=2,
    )
    result = contracts.EvidenceResult(status="ok", evidence=(evidence,))
    payload = result.to_dict()
    assert payload["has_valid_evidence"] is True and "answer" not in payload
    expect_error(lambda: contracts.EvidenceResult(status="ok"), "contain evidence")
    expect_error(lambda: contracts.EvidenceResult(status="invalid"), "status")
    expect_error(lambda: contracts.Evidence(chunk_id="x", content="x", source="Doc.docx", source_path="Doc.docx", page_number=1), "DOCX")

    content_v1 = b"price=79"
    content_v2 = b"price=99"
    with tempfile.TemporaryDirectory() as tmp:
        path = Path(tmp) / "manifest.json"
        manifest = manifest_module.Manifest.load(path)
        first = manifest.begin_generation("pricing/Pricing.pdf", manifest.content_hash(content_v1), ["c2", "c1"])
        assert first.active_generation is None and first.pending_generation == 1
        expect_error(lambda: manifest.activate(first.document_id, ["c1"]), "chunk set")
        assert manifest.active_snapshot() == {}
        activated = manifest.activate(first.document_id, ["c1", "c2"])
        assert activated.active_generation == 1 and manifest.active_snapshot() == {first.document_id: 1}
        assert manifest.begin_generation("pricing/Pricing.pdf", manifest.content_hash(content_v1), ["c1", "c2"]) is None

        second = manifest.begin_generation("pricing/Pricing.pdf", manifest.content_hash(content_v2), ["n1", "n2"])
        assert second.active_generation == 1 and second.pending_generation == 2
        assert manifest.active_snapshot() == {first.document_id: 1}
        failed = manifest.fail(second.document_id, "search_timeout")
        assert failed.state == "indexed" and failed.active_generation == 1
        assert manifest.active_snapshot() == {first.document_id: 1}

        retry = manifest.begin_generation("pricing/Pricing.pdf", manifest.content_hash(content_v2), ["n1", "n2"])
        assert retry.pending_generation == 2
        activated_v2 = manifest.activate(retry.document_id, (item for item in ["n1", "n2"]))
        assert activated_v2.active_generation == 2 and activated_v2.previous_generation == 1
        loaded = manifest_module.Manifest.load(path)
        assert loaded.active_snapshot() == {first.document_id: 2}
        assert not path.with_name(f".{path.name}.tmp").exists()

    config = azure_config()
    assert azure.load_config(config) == config
    missing = dict(config); del missing["AZURE_SEARCH_QUERY_KEY"]
    expect_error(lambda: azure.load_config(missing), "AZURE_SEARCH_QUERY_KEY")
    insecure = dict(config); insecure["AZURE_SEARCH_ENDPOINT"] = "http://search.example"
    expect_error(lambda: azure.load_config(insecure), "https")

    calls = []
    sleeps = []
    responses = [
        azure.Response(429, {"Retry-After": "2"}, b'{"error":"throttled"}'),
        azure.Response(503, {}, b"untrusted source body with admin-secret"),
        azure.Response(200, {}, b'{"value":[]}'),
    ]
    def fake_transport(method, url, headers, body, timeout):
        calls.append((method, url, dict(headers), body, timeout))
        return responses.pop(0)

    client = azure.AzureClient(fake_transport, timeout=4, max_attempts=3, sleeper=sleeps.append, jitter=lambda: 0)
    search = azure.SearchClients(config, client)
    assert search.query({"search": "pricing"}).json() == {"value": []}
    assert len(calls) == 3 and sleeps == [2.0, 2]
    assert all(call[2]["api-key"] == "query-secret" for call in calls)
    assert all("admin-secret" not in str(call[3]) for call in calls)

    mutation_calls = []
    mutation = azure.SearchClients(config, azure.AzureClient(lambda *args: mutation_calls.append(args) or azure.Response(200, {}, b"{}")))
    mutation.mutate({"value": []})
    assert mutation_calls[0][2]["api-key"] == "admin-secret"

    for status, code in ((401, "unauthorized"), (403, "unauthorized"), (400, "request_failed"), (500, "retry_exhausted")):
        failing = azure.AzureClient(lambda *args, value=status: azure.Response(value, {}, b"customer-private-body"), max_attempts=1)
        expect_error(lambda current=failing: current.request("GET", "https://service.example"), code)
    timeout = azure.AzureClient(lambda *args: (_ for _ in ()).throw(azure.AzureError("azure_timeout")), max_attempts=2, sleeper=lambda _: None)
    expect_error(lambda: timeout.request("GET", "https://service.example"), "timeout")
    invalid = azure.Response(200, {}, b"not-json")
    expect_error(invalid.json, "invalid_json")

    sensitive_values = [config[name] for name in azure.SECRET_ENV]
    redacted = azure.redact(" ".join(sensitive_values), config)
    assert all(value not in redacted for value in sensitive_values)
    assert redacted.count("[REDACTED]") == len(sensitive_values)
    for error in (
        azure.AzureError("azure_unauthorized", 401), azure.AzureError("azure_request_failed", 400),
        azure.AzureError("azure_retry_exhausted", 500), azure.AzureError("azure_timeout"),
    ):
        rendered = str(error)
        assert "customer-private-body" not in rendered and "secret" not in rendered

    markdown_units = extraction.extract("handbook/policy.md", b"# Refunds\n\nApproval required.\n\n# Travel\n\nReceipts required.")
    assert [unit.section_heading for unit in markdown_units] == ["Refunds", "Travel"]
    assert markdown_units[0].line_range == "3-3"
    html_units = extraction.extract("web/policy.html", b"<h1>Policy</h1><p>Keep records.</p>")
    assert html_units and "Keep records" in html_units[0].text
    csv_units = extraction.extract("data/prices.csv", b"sku,price\nPRO,99\nBASIC,20\n")
    assert csv_units[0].text == "sku: PRO | price: 99" and csv_units[0].line_range == "2"
    expect_error(lambda: extraction.extract("word/policy.docx", b"fake"), "python-docx")

    ocr_calls = []
    ocr_units = extraction.extract(
        "scan/policy.pdf", b"not-a-pdf",
        lambda content: ocr_calls.append(content) or [extraction.Unit("Scanned rule", page_number=1, extraction_method="document-intelligence", is_ocr=True)],
    )
    assert ocr_calls == [b"not-a-pdf"] and ocr_units[0].is_ocr

    chunks = extraction.chunk_units("doc", 2, [extraction.Unit("one two three four five", page_number=4)], target_words=3, overlap_words=1)
    assert [chunk.chunk_id for chunk in chunks] == ["doc:g2:c1", "doc:g2:c2", "doc:g2:c3"]
    assert all(chunk.page_number == 4 for chunk in chunks)

    from pptx import Presentation
    ppt = Presentation(); slide = ppt.slides.add_slide(ppt.slide_layouts[5]); slide.shapes.title.text = "Quarterly Plan"
    ppt_buffer = __import__("io").BytesIO(); ppt.save(ppt_buffer)
    ppt_units = extraction.extract("slides/plan.pptx", ppt_buffer.getvalue())
    assert ppt_units[0].slide_number == 1 and "Quarterly Plan" in ppt_units[0].text

    from openpyxl import Workbook
    workbook = Workbook(); sheet = workbook.active; sheet.title = "Pricing"; sheet.append(["sku", "price"]); sheet.append(["PRO", 99])
    xlsx_buffer = __import__("io").BytesIO(); workbook.save(xlsx_buffer)
    xlsx_units = extraction.extract("tables/pricing.xlsx", xlsx_buffer.getvalue())
    assert xlsx_units[0].sheet_name == "Pricing" and xlsx_units[0].cell_range == "A2:B2"

    with tempfile.TemporaryDirectory() as tmp:
        ingest_manifest = manifest_module.Manifest(Path(tmp) / "manifest.json")
        uploads = []; deletes = []
        service = ingestion.IngestionService(
            ingest_manifest,
            lambda texts: [[float(i + 1)] for i, _ in enumerate(texts)],
            uploads.extend,
            lambda document_id, generation: [item["chunk_id"] for item in uploads if item["document_id"] == document_id and item["generation"] == generation],
            lambda document_id, generation: deletes.append((document_id, generation)),
        )
        first_result = service.ingest("kb/policy.md", b"# Policy\n\nApproval required.", ["internal"])
        assert first_result["status"] == "indexed" and first_result["generation"] == 1
        assert uploads[0]["access_groups"] == ["internal"] and uploads[0]["content_vector"] == [1.0]
        unchanged = service.ingest("kb/policy.md", b"# Policy\n\nApproval required.", ["internal"])
        assert unchanged["status"] == "unchanged" and len(uploads) == 1
        second_result = service.ingest("kb/policy.md", b"# Policy\n\nManager approval required.", ["internal"])
        assert second_result["generation"] == 2 and deletes == [(first_result["document_id"], 1)]

        failing_manifest = manifest_module.Manifest(Path(tmp) / "failed.json")
        failing_service = ingestion.IngestionService(failing_manifest, lambda texts: [[1.0]], lambda docs: None, lambda document_id, generation: [], lambda document_id, generation: None)
        expect_error(lambda: failing_service.ingest("kb/fail.md", b"# Fail\n\nText", ["internal"]), "chunk set")
        failed_record = next(iter(failing_manifest.records.values()))
        assert failed_record.state == "failed" and failed_record.active_generation is None


if __name__ == "__main__":
    parser = ArgumentParser()
    parser.add_argument("--layer", type=int, choices=(1, 2), required=True)
    args = parser.parse_args()
    (layer_1 if args.layer == 1 else layer_2)()
    print(f"knowledge layer {args.layer}: pass")
