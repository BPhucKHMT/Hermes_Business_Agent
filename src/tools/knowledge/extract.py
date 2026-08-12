from dataclasses import dataclass
from html.parser import HTMLParser
from pathlib import Path
import csv
import io
from typing import Callable, Iterable, List, Optional, Sequence

from contracts import validate_source_path


class ExtractionError(ValueError):
    pass


@dataclass(frozen=True)
class Unit:
    text: str
    section_heading: Optional[str] = None
    page_number: Optional[int] = None
    slide_number: Optional[int] = None
    sheet_name: Optional[str] = None
    cell_range: Optional[str] = None
    line_range: Optional[str] = None
    extraction_method: str = "native"
    is_ocr: bool = False


@dataclass(frozen=True)
class Chunk:
    chunk_id: str
    ordinal: int
    content: str
    section_heading: Optional[str] = None
    page_number: Optional[int] = None
    slide_number: Optional[int] = None
    sheet_name: Optional[str] = None
    cell_range: Optional[str] = None
    line_range: Optional[str] = None
    extraction_method: str = "native"
    is_ocr: bool = False


class _HTMLText(HTMLParser):
    def __init__(self):
        HTMLParser.__init__(self)
        self.parts = []
        self.heading = None
        self._heading_tag = None

    def handle_starttag(self, tag, attrs):
        if tag in {"h1", "h2", "h3", "h4", "h5", "h6"}:
            self._heading_tag = tag
        if tag in {"p", "div", "li", "tr", "br"}:
            self.parts.append("\n")

    def handle_endtag(self, tag):
        if tag == self._heading_tag:
            self._heading_tag = None
        if tag in {"p", "div", "li", "tr"}:
            self.parts.append("\n")

    def handle_data(self, data):
        value = data.strip()
        if value:
            if self._heading_tag:
                self.heading = value
            self.parts.append(value + " ")


def _text_units(text: str, method: str = "native-text") -> List[Unit]:
    lines = text.replace("\r\n", "\n").replace("\r", "\n").split("\n")
    units = []
    heading = None
    paragraph = []
    start = 1

    def flush(end):
        if paragraph:
            units.append(Unit(" ".join(paragraph).strip(), heading, line_range="%d-%d" % (start, end), extraction_method=method))
            paragraph[:] = []

    for number, raw in enumerate(lines, 1):
        value = raw.strip()
        if value.startswith("#") and value.lstrip("#").startswith(" "):
            flush(number - 1)
            heading = value.lstrip("#").strip()
            start = number + 1
        elif not value:
            flush(number - 1)
            start = number + 1
        else:
            if not paragraph:
                start = number
            paragraph.append(value)
    flush(len(lines))
    return [unit for unit in units if unit.text]


def extract(source_path: str, content: bytes, ocr: Optional[Callable[[bytes], Sequence[Unit]]] = None, min_chars_per_page: int = 80) -> List[Unit]:
    normalized = validate_source_path(source_path)
    suffix = Path(normalized).suffix.lower()
    if suffix in {".txt", ".md"}:
        return _text_units(content.decode("utf-8-sig"), "native-%s" % suffix[1:])
    if suffix == ".html":
        parser = _HTMLText(); parser.feed(content.decode("utf-8-sig"))
        return _text_units("".join(parser.parts), "native-html")
    if suffix == ".csv":
        rows = list(csv.reader(io.StringIO(content.decode("utf-8-sig"))))
        if not rows:
            return []
        header = rows[0]
        return [Unit(" | ".join("%s: %s" % (header[i] if i < len(header) else "column_%d" % (i + 1), value) for i, value in enumerate(row)), line_range=str(number), extraction_method="native-csv") for number, row in enumerate(rows[1:], 2)]
    if suffix == ".pdf":
        try:
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(content))
            units = [Unit((page.extract_text() or "").strip(), page_number=i, extraction_method="native-pdf") for i, page in enumerate(reader.pages, 1)]
        except Exception as exc:
            if ocr is None:
                raise ExtractionError("PDF parser failed; configure Document Intelligence") from exc
            return list(ocr(content))
        nonempty = [unit for unit in units if unit.text]
        average = sum(len(unit.text) for unit in units) / max(len(units), 1)
        if not nonempty or average < min_chars_per_page:
            if ocr is None:
                raise ExtractionError("PDF requires Document Intelligence OCR")
            return list(ocr(content))
        return nonempty
    if suffix == ".pptx":
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise ExtractionError("PPTX parser unavailable; install python-pptx") from exc
        presentation = Presentation(io.BytesIO(content))
        return [Unit("\n".join(shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()), slide_number=i, extraction_method="native-pptx") for i, slide in enumerate(presentation.slides, 1)]
    if suffix == ".xlsx":
        try:
            from openpyxl import load_workbook
        except ImportError as exc:
            raise ExtractionError("XLSX parser unavailable; install openpyxl") from exc
        workbook = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        units = []
        for sheet in workbook.worksheets:
            rows = list(sheet.iter_rows(values_only=True))
            if not rows:
                continue
            headers = [str(value or "column_%d" % (i + 1)) for i, value in enumerate(rows[0])]
            for number, row in enumerate(rows[1:], 2):
                text = " | ".join("%s: %s" % (headers[i], value) for i, value in enumerate(row) if value is not None)
                if text:
                    units.append(Unit(text, sheet_name=sheet.title, cell_range="A%d:%s%d" % (number, _column_name(len(row)), number), extraction_method="native-xlsx"))
        return units
    if suffix == ".docx":
        raise ExtractionError("DOCX parser unavailable; install python-docx before claiming support")
    raise ExtractionError("unsupported source type")


def _column_name(number: int) -> str:
    result = ""
    while number:
        number, remainder = divmod(number - 1, 26)
        result = chr(65 + remainder) + result
    return result or "A"


def chunk_units(document_id: str, generation: int, units: Iterable[Unit], target_words: int = 375, overlap_words: int = 60) -> List[Chunk]:
    if target_words < 1 or overlap_words < 0 or overlap_words >= target_words:
        raise ValueError("invalid chunk limits")
    chunks = []
    ordinal = 0
    step = target_words - overlap_words
    for unit in units:
        words = unit.text.split()
        for start in range(0, len(words), step):
            piece = words[start:start + target_words]
            if not piece:
                continue
            ordinal += 1
            chunks.append(Chunk(
                chunk_id="%s:g%d:c%d" % (document_id, generation, ordinal), ordinal=ordinal,
                content=" ".join(piece), section_heading=unit.section_heading,
                page_number=unit.page_number, slide_number=unit.slide_number,
                sheet_name=unit.sheet_name, cell_range=unit.cell_range,
                line_range=unit.line_range, extraction_method=unit.extraction_method, is_ocr=unit.is_ocr,
            ))
    return chunks
