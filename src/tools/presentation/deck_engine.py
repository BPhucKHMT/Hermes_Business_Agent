"""
Hermes Executive Presentation & Visual Studio Engine (deck_engine.py)
Universal C-Level 16:9 Presentation Generator for all business intents across all Hermes profiles.
"""

from argparse import ArgumentParser
from pathlib import Path
import json
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    Presentation = None


def hex_to_rgb(hex_str: str) -> RGBColor:
    hex_str = hex_str.lstrip("#")
    return RGBColor(*(int(hex_str[i:i+2], 16) for i in (0, 2, 4)))


THEMES = {
    "dark": {
        "bg": hex_to_rgb("0B132B"),         # Deep Space Navy
        "card_bg": hex_to_rgb("1C2541"),    # Elevated Dark Slate Card
        "card_border": hex_to_rgb("3A506B"),# Slate Border
        "text_main": hex_to_rgb("FFFFFF"),  # Primary White
        "text_sub": hex_to_rgb("94A3B8"),   # Muted Slate
        "accent_cyan": hex_to_rgb("48CAE4"),
        "accent_emerald": hex_to_rgb("06D6A0"),
        "accent_gold": hex_to_rgb("FFD166"),
        "accent_coral": hex_to_rgb("EF476F"),
    },
    "light": {
        "bg": hex_to_rgb("F8FAFC"),         # Clean Crisp Slate
        "card_bg": hex_to_rgb("FFFFFF"),    # Pure White Card
        "card_border": hex_to_rgb("E2E8F0"),# Light Border
        "text_main": hex_to_rgb("0F172A"),  # Deep Navy Text
        "text_sub": hex_to_rgb("64748B"),   # Slate Text
        "accent_cyan": hex_to_rgb("0284C7"),
        "accent_emerald": hex_to_rgb("059669"),
        "accent_gold": hex_to_rgb("D97706"),
        "accent_coral": hex_to_rgb("DC2626"),
    }
}


class ExecutiveDeckBuilder:
    def __init__(self, theme_name="dark"):
        if Presentation is None:
            raise RuntimeError("python-pptx library is required to run ExecutiveDeckBuilder")
        self.theme = THEMES.get(theme_name, THEMES["dark"])
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.blank_layout = self.prs.slide_layouts[6]

    def _set_bg(self, slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = self.theme["bg"]

    def _add_header(self, slide, title: str, category: str = "EXECUTIVE BRIEFING"):
        tb_cat = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.35))
        p_c = tb_cat.text_frame.paragraphs[0]
        p_c.text = category.upper()
        p_c.font.size = Pt(11)
        p_c.font.bold = True
        p_c.font.color.rgb = self.theme["accent_cyan"]
        p_c.font.name = "Segoe UI"

        tb_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        p_t = tb_title.text_frame.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(22)
        p_t.font.bold = True
        p_t.font.color.rgb = self.theme["text_main"]
        p_t.font.name = "Segoe UI"

    def _add_card_shape(self, slide, left, top, width, height, bg_color=None, border_color=None):
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        sh.fill.solid()
        sh.fill.fore_color.rgb = bg_color or self.theme["card_bg"]
        sh.line.color.rgb = border_color or self.theme["card_border"]
        sh.line.width = Pt(1)
        return sh

    def add_hero_slide(self, title: str, subtitle: str, category: str = "EXECUTIVE BRIEFING", highlight: str = ""):
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._set_bg(slide)

        tb_cat = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(0.4))
        p_c = tb_cat.text_frame.paragraphs[0]
        p_c.text = category.upper()
        p_c.font.size = Pt(12)
        p_c.font.bold = True
        p_c.font.color.rgb = self.theme["accent_cyan"]
        p_c.font.name = "Segoe UI"

        tb_title = slide.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(11.3), Inches(2.2))
        p_t = tb_title.text_frame.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(36)
        p_t.font.bold = True
        p_t.font.color.rgb = self.theme["text_main"]
        p_t.font.name = "Segoe UI"

        if subtitle or highlight:
            self._add_card_shape(slide, Inches(1.0), Inches(4.8), Inches(11.33), Inches(1.6))
            tb_body = slide.shapes.add_textbox(Inches(1.25), Inches(4.95), Inches(10.8), Inches(1.3))
            tf = tb_body.text_frame
            tf.word_wrap = True
            if highlight:
                p_h = tf.paragraphs[0]
                p_h.text = highlight
                p_h.font.size = Pt(15)
                p_h.font.bold = True
                p_h.font.color.rgb = self.theme["accent_emerald"]
                p_h.font.name = "Segoe UI"
            if subtitle:
                p_s = tf.add_paragraph() if highlight else tf.paragraphs[0]
                p_s.text = subtitle
                p_s.font.size = Pt(12)
                p_s.font.color.rgb = self.theme["text_main"]
                p_s.font.name = "Segoe UI"
                p_s.space_before = Pt(4)

    def add_stats_slide(self, title: str, stats_list: list, category: str = "KEY PERFORMANCE INDICATORS"):
        """stats_list = [{'stat': '1.2 TY', 'label': 'TOTAL BUDGET', 'desc': 'Ngan sach tieu chuan', 'accent': 'accent_gold'}]"""
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._set_bg(slide)
        self._add_header(slide, title, category)

        count = min(len(stats_list), 4)
        if count == 0:
            return
        col_width = (Inches(11.73) - Inches(0.3) * (count - 1)) / count
        start_left = Inches(0.8)

        for i, item in enumerate(stats_list[:count]):
            left = start_left + i * (col_width + Inches(0.3))
            accent_key = item.get("accent", "accent_cyan")
            accent_color = self.theme.get(accent_key, self.theme["accent_cyan"])

            self._add_card_shape(slide, left, Inches(1.8), col_width, Inches(5.0))
            tb = slide.shapes.add_textbox(left + Inches(0.2), Inches(2.0), col_width - Inches(0.4), Inches(4.6))
            tf = tb.text_frame
            tf.word_wrap = True

            p_s = tf.paragraphs[0]
            p_s.text = str(item.get("stat", ""))
            p_s.font.size = Pt(32)
            p_s.font.bold = True
            p_s.font.color.rgb = accent_color
            p_s.font.name = "Segoe UI"

            p_l = tf.add_paragraph()
            p_l.text = str(item.get("label", "")).upper()
            p_l.font.size = Pt(13)
            p_l.font.bold = True
            p_l.font.color.rgb = self.theme["text_main"]
            p_l.font.name = "Segoe UI"
            p_l.space_before = Pt(8)

            p_d = tf.add_paragraph()
            p_d.text = str(item.get("desc", ""))
            p_d.font.size = Pt(11)
            p_d.font.color.rgb = self.theme["text_sub"]
            p_d.font.name = "Segoe UI"
            p_d.space_before = Pt(8)

    def add_quadrant_slide(self, title: str, quadrants: list, category: str = "STRATEGIC ANALYSIS"):
        """quadrants = [{'title': '1. Mat bang', 'bullets': ['Diem A', 'Diem B'], 'accent': 'accent_cyan'}] up to 4 items"""
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._set_bg(slide)
        self._add_header(slide, title, category)

        positions = [
            (Inches(0.8), Inches(1.8)),
            (Inches(6.8), Inches(1.8)),
            (Inches(0.8), Inches(4.5)),
            (Inches(6.8), Inches(4.5)),
        ]
        accents = ["accent_cyan", "accent_emerald", "accent_gold", "accent_coral"]

        for idx, item in enumerate(quadrants[:4]):
            left, top = positions[idx]
            accent_key = item.get("accent", accents[idx])
            accent_color = self.theme.get(accent_key, self.theme["accent_cyan"])

            self._add_card_shape(slide, left, top, Inches(5.7), Inches(2.4))
            tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), Inches(5.3), Inches(2.1))
            tf = tb.text_frame
            tf.word_wrap = True

            p_t = tf.paragraphs[0]
            p_t.text = item.get("title", f"Muc #{idx+1}")
            p_t.font.size = Pt(14)
            p_t.font.bold = True
            p_t.font.color.rgb = accent_color
            p_t.font.name = "Segoe UI"

            bullets = item.get("bullets", [])
            body_text = "\n".join(f"• {b}" for b in bullets) if isinstance(bullets, list) else str(bullets)
            p_b = tf.add_paragraph()
            p_b.text = body_text
            p_b.font.size = Pt(11)
            p_b.font.color.rgb = self.theme["text_main"]
            p_b.font.name = "Segoe UI"
            p_b.space_before = Pt(4)

    def add_card_grid_slide(self, title: str, cards: list, category: str = "STRATEGIC WORKSTREAMS"):
        """3 vertical cards layout"""
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._set_bg(slide)
        self._add_header(slide, title, category)

        count = min(len(cards), 3)
        col_width = Inches(3.75)
        gap = Inches(0.24)
        start_left = Inches(0.8)
        accents = ["accent_cyan", "accent_emerald", "accent_gold"]

        for idx, card in enumerate(cards[:count]):
            left = start_left + idx * (col_width + gap)
            accent_key = card.get("accent", accents[idx % 3])
            accent_color = self.theme.get(accent_key, self.theme["accent_cyan"])

            self._add_card_shape(slide, left, Inches(1.8), col_width, Inches(5.0))
            tb = slide.shapes.add_textbox(left + Inches(0.2), Inches(1.95), col_width - Inches(0.4), Inches(4.7))
            tf = tb.text_frame
            tf.word_wrap = True

            p_t = tf.paragraphs[0]
            p_t.text = card.get("title", f"Card #{idx+1}")
            p_t.font.size = Pt(15)
            p_t.font.bold = True
            p_t.font.color.rgb = accent_color
            p_t.font.name = "Segoe UI"

            bullets = card.get("bullets", [])
            body = "\n\n".join(f"• {b}" for b in bullets) if isinstance(bullets, list) else str(bullets)
            p_b = tf.add_paragraph()
            p_b.text = body
            p_b.font.size = Pt(11.5)
            p_b.font.color.rgb = self.theme["text_main"]
            p_b.font.name = "Segoe UI"
            p_b.space_before = Pt(8)

    def add_timeline_slide(self, title: str, phases: list, category: str = "EXECUTION ROADMAP"):
        """3 or 4 horizontal phase blocks"""
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._set_bg(slide)
        self._add_header(slide, title, category)

        count = min(len(phases), 3)
        col_width = Inches(3.75)
        gap = Inches(0.24)
        start_left = Inches(0.8)
        accents = ["accent_cyan", "accent_emerald", "accent_gold"]

        for idx, phase in enumerate(phases[:count]):
            left = start_left + idx * (col_width + gap)
            accent_color = self.theme[accents[idx % 3]]

            self._add_card_shape(slide, left, Inches(1.8), col_width, Inches(5.0))
            tb = slide.shapes.add_textbox(left + Inches(0.2), Inches(1.95), col_width - Inches(0.4), Inches(4.7))
            tf = tb.text_frame
            tf.word_wrap = True

            p_ph = tf.paragraphs[0]
            p_ph.text = f"PHASE {idx+1}: {phase.get('timeframe', '').upper()}"
            p_ph.font.size = Pt(11)
            p_ph.font.bold = True
            p_ph.font.color.rgb = accent_color
            p_ph.font.name = "Segoe UI"

            p_t = tf.add_paragraph()
            p_t.text = phase.get("title", f"Giai đoạn {idx+1}")
            p_t.font.size = Pt(15)
            p_t.font.bold = True
            p_t.font.color.rgb = self.theme["text_main"]
            p_t.font.name = "Segoe UI"
            p_t.space_before = Pt(4)

            tasks = phase.get("tasks", [])
            body = "\n".join(f"✓ {t}" for t in tasks) if isinstance(tasks, list) else str(tasks)
            p_b = tf.add_paragraph()
            p_b.text = body
            p_b.font.size = Pt(11)
            p_b.font.color.rgb = self.theme["text_sub"]
            p_b.font.name = "Segoe UI"
            p_b.space_before = Pt(8)

    def save(self, output_path: Path) -> Path:
        out = Path(output_path)
        out.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(out))
        return out


def render_from_spec(spec_data: dict, output_path: Path, theme: str = "dark") -> Path:
    builder = ExecutiveDeckBuilder(theme_name=theme)
    for slide in spec_data.get("slides", []):
        layout = slide.get("layout", "hero")
        title = slide.get("title", "")
        cat = slide.get("category", spec_data.get("category", "EXECUTIVE BRIEFING"))

        if layout == "hero":
            builder.add_hero_slide(title, slide.get("subtitle", ""), cat, highlight=slide.get("highlight", ""))
        elif layout == "stats":
            builder.add_stats_slide(title, slide.get("stats", []), cat)
        elif layout == "quadrant":
            builder.add_quadrant_slide(title, slide.get("quadrants", []), cat)
        elif layout == "cards":
            builder.add_card_grid_slide(title, slide.get("cards", []), cat)
        elif layout == "timeline":
            builder.add_timeline_slide(title, slide.get("phases", []), cat)

    return builder.save(output_path)


def main():
    parser = ArgumentParser(description="Render C-level executive presentation deck from JSON spec")
    parser.add_argument("spec", type=Path, help="JSON specification file")
    parser.add_argument("output", type=Path, help="Output .pptx filepath")
    parser.add_argument("--theme", choices=("dark", "light"), default="dark")
    args = parser.parse_args()

    spec_data = json.loads(args.spec.read_text(encoding="utf-8"))
    res = render_from_spec(spec_data, args.output, theme=args.theme)
    print(f"Executive presentation successfully built at: {res.resolve()}")


if __name__ == "__main__":
    main()
