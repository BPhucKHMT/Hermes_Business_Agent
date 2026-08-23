import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

spec = json.load(open('workspaces/protein-bar/week-3-checklist.json', encoding='utf-8'))
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

def rgb(h): return RGBColor.from_string(h)
def box(slide, left, top, width, height, fill):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    sh.fill.solid(); sh.fill.fore_color.rgb = rgb(fill); sh.line.fill.background(); return sh
def text(slide, x, y, w, h, value, size, color='12372A', bold=False, align=None):
    tb=slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf=tb.text_frame; tf.clear(); tf.word_wrap=True
    p=tf.paragraphs[0]; p.text=value; p.font.name='Aptos'; p.font.size=Pt(size); p.font.bold=bold; p.font.color.rgb=rgb(color)
    if align: p.alignment=align
    return tb
for s in spec['slides']:
    slide=prs.slides.add_slide(blank); bg=s.get('background','F6F3EC'); box(slide,0,0,13.333,7.5,bg)
    title=s.get('title',''); title_color='FFFFFF' if bg=='12372A' else '12372A'
    text(slide,0.7,0.55,11.9,0.75,title,28,title_color,True)
    if s['layout']=='title':
        text(slide,0.75,2.25,11.8,1.5,s.get('subtitle',''),22,'FFFFFF',False)
        box(slide,0.75,5.55,2.0,0.08,'E6B566')
    else:
        y=1.55
        for b in s.get('bullets',[]):
            if isinstance(b,str): b={'text':b}
            text(slide,0.9,y,11.5,0.62,'• '+b['text'],18,b.get('color','12372A'),b.get('bold',False)); y+=0.73
    footer=s.get('footer','')
    text(slide,0.7,7.08,11.9,0.22,footer,9,'FFFFFF' if bg=='12372A' else '567064')
prs.save('workspaces/protein-bar/Protein_Bar_Week_3_Checklist.pptx')
print('created',len(prs.slides))
