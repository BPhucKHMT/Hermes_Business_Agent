from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

def hex_to_rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(*(int(h[i:i+2], 16) for i in (0, 2, 4)))

# Luxury Executive Palette (Deep Slate & Vibrant Accents)
BG_DARK = hex_to_rgb("0B132B")        # Deep Space Navy
CARD_BG = hex_to_rgb("1C2541")        # Elevated Slate Card
CARD_BORDER = hex_to_rgb("3A506B")    # Subtle Slate Border
TEXT_WHITE = hex_to_rgb("FFFFFF")     # Crisp White
TEXT_MUTED = hex_to_rgb("94A3B8")     # Subtext Muted Slate
CYAN_ACCENT = hex_to_rgb("48CAE4")    # Neon Cyan
EMERALD_ACCENT = hex_to_rgb("06D6A0") # Fresh Emerald
GOLD_ACCENT = hex_to_rgb("FFD166")    # Rich Gold
CORAL_ACCENT = hex_to_rgb("EF476F")   # Alert Coral

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

def set_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_header(slide, title: str, category: str = "PROTEIN BAR THẢO ĐIỀN — WEEK 3 OPERATIONS"):
    # Category Tag
    tb_cat = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.35))
    p_c = tb_cat.text_frame.paragraphs[0]
    p_c.text = category.upper()
    p_c.font.size = Pt(11)
    p_c.font.bold = True
    p_c.font.color.rgb = CYAN_ACCENT
    p_c.font.name = "Segoe UI"

    # Main Title
    tb_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
    p_t = tb_title.text_frame.paragraphs[0]
    p_t.text = title
    p_t.font.size = Pt(22)
    p_t.font.bold = True
    p_t.font.color.rgb = TEXT_WHITE
    p_t.font.name = "Segoe UI"

def add_card(slide, left, top, width, height, title: str, body: str, accent=CYAN_ACCENT, stat: str = ""):
    # Card Background Shape with Border
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = CARD_BG
    sh.line.color.rgb = CARD_BORDER
    sh.line.width = Pt(1)

    tb = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.2), width - Inches(0.5), height - Inches(0.4))
    tf = tb.text_frame
    tf.word_wrap = True

    if stat:
        p_stat = tf.paragraphs[0]
        p_stat.text = stat
        p_stat.font.size = Pt(32)
        p_stat.font.bold = True
        p_stat.font.color.rgb = accent
        p_stat.font.name = "Segoe UI"
        
        p_t = tf.add_paragraph()
        p_t.text = title
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_WHITE
        p_t.font.name = "Segoe UI"
        p_t.space_before = Pt(4)
    else:
        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(15)
        p_t.font.bold = True
        p_t.font.color.rgb = accent
        p_t.font.name = "Segoe UI"

    p_b = tf.add_paragraph()
    p_b.text = body
    p_b.font.size = Pt(11.5)
    p_b.font.color.rgb = TEXT_WHITE
    p_b.font.name = "Segoe UI"
    p_b.space_before = Pt(6)

# ==========================================
# SLIDE 1: Cover Hero
# ==========================================
s1 = prs.slides.add_slide(blank)
set_bg(s1, BG_DARK)

tb1 = s1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(0.4))
p1 = tb1.text_frame.paragraphs[0]
p1.text = "EXECUTIVE OPERATIONAL CHECKLIST"
p1.font.size = Pt(12)
p1.font.bold = True
p1.font.color.rgb = CYAN_ACCENT

tb_main = s1.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(11.3), Inches(2.0))
pm = tb_main.text_frame.paragraphs[0]
pm.text = "Kế Hoạch Triển Khai Tuần 3\n(24/08 – 30/08/2026)"
pm.font.size = Pt(36)
pm.font.bold = True
pm.font.color.rgb = TEXT_WHITE

add_card(s1, Inches(1.0), Inches(4.8), Inches(11.33), Inches(1.6), 
         "Mục Tiêu Trọng Tâm Tuần 3", 
         "Chốt địa điểm mặt bằng Thảo Điền (kiểm tra ngập nước) • Hoàn thiện hồ sơ ĐKKD • Tasting 5 SKU RTD đầu tiên • Khóa dự toán COGS menu.", 
         EMERALD_ACCENT)

# ==========================================
# SLIDE 2: KPI Numbers & Hard Gates
# ==========================================
s2 = prs.slides.add_slide(blank)
set_bg(s2, BG_DARK)
add_header(s2, "Chỉ Số Cốt Lõi & Khung Thời Gian Mở Quán")

add_card(s2, Inches(0.8), Inches(1.8), Inches(2.75), Inches(5.0), 
         "TARGET OPENING", "Mục tiêu vận hành thử 05/12, khai trương chính thức trước 08/12.", CYAN_ACCENT, stat="05/12")

add_card(s2, Inches(3.8), Inches(1.8), Inches(2.75), Inches(5.0), 
         "TOTAL BUDGET", "Ngân sách tiêu chuẩn 1.2 tỷ VND (Dự phòng rủi ro 1.44 tỷ VND).", GOLD_ACCENT, stat="1.2 TỶ")

add_card(s2, Inches(6.8), Inches(1.8), Inches(2.75), Inches(5.0), 
         "TASK TRACKING", "Tổng cộng 143 đầu việc chia đều trong 16 tuần triển khai liên tục.", EMERALD_ACCENT, stat="143 TASKS")

add_card(s2, Inches(9.8), Inches(1.8), Inches(2.75), Inches(5.0), 
         "HARD GATES", "3 Cổng pháp lý bắt buộc: ĐKKD (Tuần 4), ATVSTP (Tuần 8), PCCC (Tuần 8).", CORAL_ACCENT, stat="4 GATES")

# ==========================================
# SLIDE 3: 4 Trọng Tâm Chi Tiết Tuần 3 (2x2 Grid)
# ==========================================
s3 = prs.slides.add_slide(blank)
set_bg(s3, BG_DARK)
add_header(s3, "Chi Tiết 4 Nhóm Công Việc Tuần 3 (24/08 – 30/08)")

add_card(s3, Inches(0.8), Inches(1.8), Inches(5.7), Inches(2.4),
         "1. Khảo Sát & Chốt Mặt Bằng (Thảo Điền)",
         "• Kiểm tra lịch sử ngập triều cường tại vị trí thuê.\n• Đàm phán thời gian fit-out miễn phí 30–45 ngày.\n• Xác minh giấy tờ quyền sở hữu hợp pháp của chủ nhà.",
         CYAN_ACCENT)

add_card(s3, Inches(6.8), Inches(1.8), Inches(5.7), Inches(2.4),
         "2. Pháp Lý & Thành Lập Doanh Nghiệp",
         "• Nộp hồ sơ thành lập Công ty TNHH 2TV tại Sở KH&ĐT.\n• Chuẩn bị đăng ký ngành nghề F&B, bán lẻ thực phẩm bổ sung.\n• Chuẩn bị hồ sơ cơ sở đủ điều kiện ATVSTP.",
         EMERALD_ACCENT)

add_card(s3, Inches(0.8), Inches(4.5), Inches(5.7), Inches(2.4),
         "3. Thử Mẫu (Tasting) & Chọn Nhà Cung Cấp",
         "• Tasting 5 SKU RTD từ Wana Beverage & Interfresh Củ Chi.\n• Yêu cầu bản Tự công bố sản phẩm và chứng nhận ISO/HACCP.\n• Đàm phán hạn mức thanh toán công nợ 15-30 ngày.",
         GOLD_ACCENT)

add_card(s3, Inches(6.8), Inches(4.5), Inches(5.7), Inches(2.4),
         "4. Menu Engineering & Định Mức COGS",
         "• Khóa công thức chuẩn cho 8 món Signature Protein Shake.\n• Thiết lập định mức COGS mục tiêu: 28% – 32% giá bán lẻ.\n• Chốt danh mục thiết bị quầy bar (Máy Espresso, Blender công nghiệp).",
         CORAL_ACCENT)

out_file = Path("src/workspaces/protein-bar/Protein_Bar_Week_3_Checklist.pptx")
prs.save(str(out_file))
print("Successfully generated luxury deck at:", out_file.resolve())
