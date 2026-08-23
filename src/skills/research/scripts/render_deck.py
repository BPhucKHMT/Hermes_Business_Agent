"""
Executive Deck Generator for Hermes Research Dossiers.
Converts canonical dossier.json into a 16:9 widescreen, C-Level PowerPoint Presentation.
"""

from argparse import ArgumentParser
from pathlib import Path
import json
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    Presentation = None


def hex_to_rgb(hex_str: str) -> RGBColor:
    hex_str = hex_str.lstrip("#")
    return RGBColor(*(int(hex_str[i:i+2], 16) for i in (0, 2, 4)))


# Theme Colors: Executive Dark Modern
BG_COLOR = hex_to_rgb("0F172A")       # Deep Slate / Navy
CARD_BG = hex_to_rgb("1E293B")        # Card Dark Slate
CARD_BORDER = hex_to_rgb("334155")    # Border Slate
TEXT_WHITE = hex_to_rgb("FFFFFF")     # Primary White
TEXT_MUTED = hex_to_rgb("94A3B8")     # Muted Slate
ACCENT_CYAN = hex_to_rgb("38BDF8")    # Highlight Blue / Cyan
ACCENT_EMERALD = hex_to_rgb("34D399") # Success Emerald
ACCENT_AMBER = hex_to_rgb("FBBF24")   # Warning Amber


def set_slide_background(slide, color=BG_COLOR):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_header(slide, title_text: str, category_text: str = "EXECUTIVE RESEARCH DOSSIER"):
    # Category badge / subtitle
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(11.5), Inches(0.35))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.size = Pt(11)
    p_cat.font.bold = True
    p_cat.font.color.rgb = ACCENT_CYAN
    p_cat.font.name = "Arial"

    # Slide Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.5), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_WHITE
    p_title.font.name = "Arial"


def add_card(slide, left, top, width, height, title: str, body: str, accent_color=ACCENT_CYAN, subtext: str = ""):
    # Card Background Shape
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = CARD_BORDER
    shape.line.width = Pt(1)

    # Card Content Text
    tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
    tf = tb.text_frame
    tf.word_wrap = True

    # Card Title
    p_t = tf.paragraphs[0]
    p_t.text = title
    p_t.font.size = Pt(15)
    p_t.font.bold = True
    p_t.font.color.rgb = accent_color
    p_t.font.name = "Arial"

    # Card Body
    p_b = tf.add_paragraph()
    p_b.text = body
    p_b.font.size = Pt(12)
    p_b.font.color.rgb = TEXT_WHITE
    p_b.font.name = "Arial"
    p_b.space_before = Pt(8)

    # Optional Subtext / Meta
    if subtext:
        p_s = tf.add_paragraph()
        p_s.text = subtext
        p_s.font.size = Pt(10)
        p_s.font.color.rgb = TEXT_MUTED
        p_s.font.name = "Arial"
        p_s.space_before = Pt(6)


def create_deck_from_dossier(dossier: dict, output_path: Path) -> Path:
    if Presentation is None:
        raise RuntimeError("python-pptx is required to generate presentation decks")

    prs = Presentation()
    # Configure 16:9 Widescreen (13.33 x 7.5 inches)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # ==========================================
    # SLIDE 1: Title / Hero Slide
    # ==========================================
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_background(s1, BG_COLOR)

    # Hero Badge
    badge_box = s1.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(10), Inches(0.4))
    p_badge = badge_box.text_frame.paragraphs[0]
    p_badge.text = "STRATEGIC INTELLIGENCE & RESEARCH REPORT"
    p_badge.font.size = Pt(12)
    p_badge.font.bold = True
    p_badge.font.color.rgb = ACCENT_CYAN

    # Main Question / Title
    title_box = s1.shapes.add_textbox(Inches(1.2), Inches(2.3), Inches(11.0), Inches(2.2))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = dossier.get("question", "Báo Cáo Nghiên Cứu")
    p1.font.size = Pt(32)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE

    # Subtitle / Scope Card
    scope_text = f"Phạm vi: {dossier.get('scope', 'Toàn diện')} | Mode: {dossier.get('mode', 'Session')} | Ngày: {dossier.get('created_at', 'N/A')[:10]}"
    add_card(s1, Inches(1.2), Inches(4.8), Inches(10.9), Inches(1.5), "Mục Tiêu & Phạm Vi Nghiên Cứu", scope_text, ACCENT_EMERALD)

    # ==========================================
    # SLIDE 2: Executive Answer & Core Takeaways
    # ==========================================
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_background(s2, BG_COLOR)
    add_header(s2, "Kết Luận Điều Hành & Đánh Giá Cốt Lõi")

    exec_answer = dossier.get("executive_answer", "Không có tóm tắt điều hành.")
    add_card(s2, Inches(0.8), Inches(1.8), Inches(11.73), Inches(2.2), "Executive Summary", exec_answer, ACCENT_CYAN)

    method_text = dossier.get("method", "Nghiên cứu đa nguồn.")
    contra_text = "; ".join(map(str, dossier.get("contradictions", []))) or "Không phát hiện mâu thuẫn trọng yếu giữa các nguồn chính."
    add_card(s2, Inches(0.8), Inches(4.3), Inches(5.7), Inches(2.5), "Phương Pháp Nghiên Cứu", method_text, ACCENT_EMERALD)
    add_card(s2, Inches(6.8), Inches(4.3), Inches(5.7), Inches(2.5), "Mâu Thuẫn & Lưu Ý Thị Trường", contra_text, ACCENT_AMBER)

    # ==========================================
    # SLIDE 3: Key Findings & Claims (Grid)
    # ==========================================
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_background(s3, BG_COLOR)
    add_header(s3, "Phát Hiện Trọng Tâm & Dữ Liệu Thực Chứng")

    claims = dossier.get("claims", [])[:3]
    col_width = Inches(3.75)
    gap = Inches(0.24)
    start_left = Inches(0.8)

    colors = [ACCENT_CYAN, ACCENT_EMERALD, ACCENT_AMBER]
    for idx, claim in enumerate(claims):
        left = start_left + idx * (col_width + gap)
        c_type = str(claim.get("type", "Fact")).upper()
        c_text = claim.get("text", "")
        rationale = f"Độ tin cậy: {claim.get('confidence', 'N/A')} ({claim.get('confidence_rationale', '')})"
        add_card(s3, left, Inches(1.8), col_width, Inches(5.0), f"[{c_type}] Luận Điểm #{idx+1}", c_text, colors[idx % len(colors)], rationale)

    # ==========================================
    # SLIDE 4: Strategic Recommendations & Gaps
    # ==========================================
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_background(s4, BG_COLOR)
    add_header(s4, "Khoảng Trống Thông Tin & Đề Xuất Hành Động")

    gaps = dossier.get("gaps", [])
    gaps_text = "\n".join(f"• {g}" for g in gaps) if gaps else "Đã bao phủ đầy đủ các khía cạnh thông tin cốt lõi."
    add_card(s4, Inches(0.8), Inches(1.8), Inches(5.7), Inches(5.0), "Khoảng Trống Cần Khảo Sát Thêm", gaps_text, ACCENT_AMBER)

    limits = dossier.get("limitations", [])
    limits_text = "\n".join(f"• {l}" for l in limits) if limits else "Áp dụng theo điều kiện thị trường công khai tiêu chuẩn."
    recs_text = f"{limits_text}\n\nĐề xuất: Triển khai kiểm chứng thực tế và cập nhật dữ liệu định kỳ."
    add_card(s4, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0), "Giới Hạn & Đề Xuất Bước Tiếp Theo", recs_text, ACCENT_EMERALD)

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


def main():
    parser = ArgumentParser(description="Render research dossier as executive PowerPoint presentation")
    parser.add_argument("dossier", type=Path, help="Path to dossier.json")
    parser.add_argument("output", type=Path, help="Path to output .pptx file")
    args = parser.parse_args()

    data = json.loads(args.dossier.read_text(encoding="utf-8"))
    out = create_deck_from_dossier(data, args.output)
    print(f"Presentation deck successfully created at: {out.resolve()}")


if __name__ == "__main__":
    main()
